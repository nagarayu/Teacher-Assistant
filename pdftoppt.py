import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from tkinter import Tk
from tkinter.filedialog import askopenfilename

def pdf_to_images(pdf_path):
    """Convert PDF to a list of images."""
    images = []
    pdf_document = fitz.open(pdf_path)
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        pix = page.get_pixmap()
        img_path = f"temp_page_{page_num}.png"
        pix.save(img_path)
        images.append(img_path)
    pdf_document.close()
    return images

def create_ppt(images, ppt_path):
    """Create a PowerPoint presentation from images."""
    prs = Presentation()
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        slide.shapes.add_picture(image, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    prs.save(ppt_path)

def main():
    # Hide the root Tkinter window
    Tk().withdraw()
    
    # Prompt user to select a PDF file
    pdf_path = askopenfilename(title="Select a PDF file", filetypes=[("PDF files", "*.pdf")])
    
    if not pdf_path:
        print("No file selected. Exiting.")
        return
    
    # Create output directory if it doesn't exist
    output_dir = "Created ppt"
    os.makedirs(output_dir, exist_ok=True)
    
    # Define the output PPT file path
    ppt_filename = os.path.splitext(os.path.basename(pdf_path))[0] + ".pptx"
    ppt_path = os.path.join(output_dir, ppt_filename)
    
    # Convert PDF to images and create PPT
    images = pdf_to_images(pdf_path)
    create_ppt(images, ppt_path)
    
    print(f"Presentation created successfully: {ppt_path}")

    # Clean up temporary image files
    for img in images:
        os.remove(img)

if __name__ == "__main__":
    main()
